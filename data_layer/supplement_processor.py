"""Convert supplement files (PDF, Word, Excel, CSV, PPTX, ZIP) to extractable text.

PDF → markdown via marker-pdf (primary) or Docling (fallback)
Word (.docx) → text via python-docx
Word (.doc) → text via OLE2 piece table parsing (olefile)
PowerPoint (.pptx) → text via python-pptx
Excel (.xlsx/.xls) → markdown tables via pandas + openpyxl
CSV/TSV → markdown tables via pandas
ZIP → extract to temp dir, process contained files

Output is either appended to the paper's main text or stored separately
for Pass 3 supplement extraction.

Usage:
    from data_layer.supplement_processor import process_supplements
    text = process_supplements("data/PMC76/PMC7612819_supp/")
"""

from __future__ import annotations

import csv
import logging
import re
import struct
import tempfile
import zipfile
from io import StringIO
from pathlib import Path

logger = logging.getLogger(__name__)

# Supplement text budget
SUPPLEMENT_TEXT_BUDGET = 40_000

# Cached PDF converter (loaded once, reused)
_pdf_converter = None
_pdf_converter_type = None  # "marker" or "docling"


def process_supplements(supp_dir: str | Path, skip_pdf: bool = False) -> str:
    """Process all supplement files in a directory into combined text.

    Args:
        supp_dir: Path to supplement directory.
        skip_pdf: If True, skip PDF files (for fast first-pass processing).

    Returns markdown text up to SUPPLEMENT_TEXT_BUDGET chars.
    """
    supp_dir = Path(supp_dir)
    if not supp_dir.exists():
        return ""

    # Collect files to process (expand ZIPs first)
    files_to_process = _collect_files(supp_dir)

    parts: list[tuple[str, str]] = []  # (filename, text)

    for fpath in files_to_process:
        ext = fpath.suffix.lower()

        # Skip very large files (>50MB) to avoid memory issues
        try:
            if fpath.stat().st_size > 50 * 1024 * 1024:
                logger.debug("Skipping oversized file: %s", fpath)
                continue
        except OSError:
            continue

        try:
            if ext == ".pdf":
                if skip_pdf:
                    continue
                text = _process_pdf(fpath)
            elif ext == ".docx":
                text = _process_docx(fpath)
            elif ext == ".doc":
                text = _process_doc_legacy(fpath)
            elif ext in (".xlsx", ".xls"):
                text = _process_excel(fpath)
            elif ext in (".csv", ".tsv"):
                text = _process_csv(fpath, delimiter="\t" if ext == ".tsv" else ",")
            elif ext == ".pptx":
                text = _process_pptx(fpath)
            elif ext == ".ppt":
                text = _process_ppt_legacy(fpath)
            elif ext == ".txt":
                text = fpath.read_text(errors="replace")[:SUPPLEMENT_TEXT_BUDGET]
            else:
                continue
        except Exception as e:
            logger.warning("Failed to process %s: %s", fpath, e)
            continue

        if text and text.strip():
            parts.append((fpath.name, text.strip()))

    if not parts:
        return ""

    # Combine with headers, respecting budget
    combined: list[str] = []
    total = 0
    for filename, text in parts:
        header = f"\n## Supplement: {filename}\n\n"
        block = header + text + "\n"
        if total + len(block) > SUPPLEMENT_TEXT_BUDGET:
            available = SUPPLEMENT_TEXT_BUDGET - total
            if available > 200:
                combined.append(block[:available - 50] + "\n[... truncated ...]\n")
            break
        combined.append(block)
        total += len(block)

    return "".join(combined)


def _collect_files(supp_dir: Path) -> list[Path]:
    """Collect processable files, extracting ZIPs into a temp directory.

    Returns a sorted list of Path objects (original files + extracted ZIP contents).
    The caller does NOT need to clean up temp dirs — they persist for the process lifetime
    which is fine since supplement processing is a batch operation.
    """
    files: list[Path] = []
    zip_files: list[Path] = []

    for fpath in sorted(supp_dir.iterdir()):
        if fpath.is_dir():
            continue
        if fpath.suffix.lower() == ".zip":
            zip_files.append(fpath)
        else:
            files.append(fpath)

    # Extract ZIP contents to a temp dir
    for zpath in zip_files:
        try:
            tmp_dir = Path(tempfile.mkdtemp(prefix="supp_zip_"))
            with zipfile.ZipFile(zpath) as zf:
                # Filter out directories, __MACOSX, and hidden files
                for name in zf.namelist():
                    if name.endswith("/") or "/__MACOSX" in name or name.startswith("__MACOSX"):
                        continue
                    basename = Path(name).name
                    if basename.startswith("."):
                        continue
                    try:
                        zf.extract(name, tmp_dir)
                        extracted = tmp_dir / name
                        if extracted.is_file():
                            files.append(extracted)
                    except Exception as e:
                        logger.debug("Failed to extract %s from %s: %s", name, zpath, e)
            logger.info("Extracted %s from ZIP: %s", len(zf.namelist()), zpath.name)
        except zipfile.BadZipFile:
            logger.warning("Bad ZIP file: %s", zpath)
        except Exception as e:
            logger.warning("Failed to open ZIP %s: %s", zpath, e)

    return sorted(files, key=lambda p: p.name)


def _get_pdf_converter():
    """Get or create a cached PDF converter (marker-pdf or Docling)."""
    global _pdf_converter, _pdf_converter_type

    if _pdf_converter is not None:
        return _pdf_converter, _pdf_converter_type

    # Try marker-pdf first
    try:
        from marker.converters.pdf import PdfConverter
        from marker.models import create_model_dict
        logger.info("Loading marker-pdf models (one-time)...")
        models = create_model_dict()
        _pdf_converter = PdfConverter(artifact_dict=models)
        _pdf_converter_type = "marker"
        logger.info("marker-pdf models loaded successfully")
        return _pdf_converter, _pdf_converter_type
    except Exception as e:
        logger.debug("marker-pdf not available: %s", e)

    # Fallback: try Docling
    try:
        from docling.document_converter import DocumentConverter
        logger.info("Loading Docling converter (one-time)...")
        _pdf_converter = DocumentConverter()
        _pdf_converter_type = "docling"
        logger.info("Docling converter loaded successfully")
        return _pdf_converter, _pdf_converter_type
    except Exception as e:
        logger.debug("Docling not available: %s", e)

    return None, None


def _process_pdf(pdf_path: Path) -> str:
    """Convert PDF to text using cached converter."""
    converter, ctype = _get_pdf_converter()
    if converter is None:
        logger.warning("No PDF converter available, skipping %s", pdf_path)
        return ""

    try:
        if ctype == "marker":
            rendered = converter(str(pdf_path))
            text = rendered.markdown
        else:  # docling
            result = converter.convert(str(pdf_path))
            text = result.document.export_to_markdown()

        if text and len(text.strip()) > 100:
            return text
    except Exception as e:
        logger.warning("PDF conversion failed for %s: %s", pdf_path, e)

    return ""


def _process_docx(docx_path: Path) -> str:
    """Convert Word .docx to text via python-docx."""
    try:
        import docx
    except ImportError:
        logger.warning("python-docx not installed, skipping %s", docx_path)
        return ""

    try:
        doc = docx.Document(str(docx_path))
        parts: list[str] = []

        for para in doc.paragraphs:
            text = para.text.strip()
            if text:
                # Detect headings
                if para.style and para.style.name and "Heading" in para.style.name:
                    level = 1
                    try:
                        level = int(para.style.name.replace("Heading", "").strip())
                    except ValueError:
                        pass
                    parts.append(f"{'#' * (level + 1)} {text}")
                else:
                    parts.append(text)

        # Also extract tables
        for table in doc.tables:
            md_table = _docx_table_to_markdown(table)
            if md_table:
                parts.append(md_table)

        return "\n\n".join(parts)

    except Exception as e:
        logger.warning("Failed to read docx %s: %s", docx_path, e)
        return ""


def _process_doc_legacy(doc_path: Path) -> str:
    """Extract text from legacy .doc (OLE2 binary Word) files via piece table parsing.

    Parses the FIB (File Information Block) and CLX (complex) structures
    to locate text pieces in the WordDocument stream.  Requires ``olefile``.
    """
    try:
        import olefile
    except ImportError:
        logger.warning("olefile not installed, skipping legacy .doc %s", doc_path)
        return ""

    try:
        ole = olefile.OleFileIO(str(doc_path))
    except Exception as e:
        logger.warning("Cannot open OLE file %s: %s", doc_path, e)
        return ""

    try:
        if not ole.exists("WordDocument"):
            return ""

        word_doc = ole.openstream("WordDocument").read()
        if len(word_doc) < 0x01AA:
            return ""

        # FIB header fields
        flags = struct.unpack_from("<H", word_doc, 0x000A)[0]
        ccp_text = struct.unpack_from("<I", word_doc, 0x004C)[0]
        if ccp_text == 0:
            return ""

        # Table stream: 1Table if bit 9 set, else 0Table
        table_name = "1Table" if (flags & 0x0200) else "0Table"
        if not ole.exists(table_name):
            table_name = "0Table" if table_name == "1Table" else "1Table"
            if not ole.exists(table_name):
                return ""

        table_data = ole.openstream(table_name).read()

        # CLX offset from FIB
        fc_clx = struct.unpack_from("<I", word_doc, 0x01A2)[0]
        lcb_clx = struct.unpack_from("<I", word_doc, 0x01A6)[0]
        if fc_clx == 0 or lcb_clx == 0 or fc_clx + lcb_clx > len(table_data):
            return ""

        clx_data = table_data[fc_clx : fc_clx + lcb_clx]

        # Walk CLX: skip Prc (0x01) entries, find Pcdt (0x02)
        pos = 0
        plc_pcd: bytes | None = None
        while pos < len(clx_data):
            clx_type = clx_data[pos]
            if clx_type == 0x02:
                pos += 1
                pcdt_size = struct.unpack_from("<I", clx_data, pos)[0]
                pos += 4
                plc_pcd = clx_data[pos : pos + pcdt_size]
                break
            elif clx_type == 0x01:
                pos += 1
                prc_size = struct.unpack_from("<H", clx_data, pos)[0]
                pos += 2 + prc_size
            else:
                pos += 1

        if plc_pcd is None or len(plc_pcd) < 16:
            return ""

        # PlcPcd: (n+1) CPs (4 bytes each) then n PCDs (8 bytes each)
        n = (len(plc_pcd) - 4) // 12
        if n <= 0:
            return ""

        cps = [struct.unpack_from("<I", plc_pcd, i * 4)[0] for i in range(n + 1)]
        pcd_offset = (n + 1) * 4

        text_parts: list[str] = []
        for i in range(n):
            char_count = cps[i + 1] - cps[i]
            if char_count <= 0:
                continue

            pcd_start = pcd_offset + i * 8
            if pcd_start + 8 > len(plc_pcd):
                break
            fc = struct.unpack_from("<I", plc_pcd, pcd_start + 2)[0]

            # Bit 30: compressed (CP1252) vs UTF-16LE
            is_compressed = bool(fc & 0x40000000)
            fc_real = fc & 0x3FFFFFFF

            if is_compressed:
                byte_off = fc_real // 2
                raw = word_doc[byte_off : byte_off + char_count]
                text_parts.append(raw.decode("cp1252", errors="replace"))
            else:
                byte_off = fc_real
                raw = word_doc[byte_off : byte_off + char_count * 2]
                text_parts.append(raw.decode("utf-16-le", errors="replace"))

        full_text = "".join(text_parts)
        if len(full_text) > ccp_text:
            full_text = full_text[:ccp_text]

        # Clean Word special characters
        full_text = full_text.replace("\r", "\n")
        full_text = full_text.replace("\x07", "\t")
        full_text = re.sub(r"[\x00-\x06\x08\x0b\x0e-\x1f]", "", full_text)
        full_text = re.sub(r"\n{3,}", "\n\n", full_text)

        return full_text.strip()

    except Exception as e:
        logger.warning("Failed to parse legacy .doc %s: %s", doc_path, e)
        return ""
    finally:
        ole.close()


def _docx_table_to_markdown(table) -> str:
    """Convert a python-docx Table to markdown."""
    rows: list[list[str]] = []
    for row in table.rows:
        cells = [cell.text.strip().replace("|", "\\|") for cell in row.cells]
        rows.append(cells)

    if not rows:
        return ""

    max_cols = max(len(r) for r in rows)
    for r in rows:
        while len(r) < max_cols:
            r.append("")

    lines: list[str] = []
    lines.append("| " + " | ".join(rows[0]) + " |")
    lines.append("| " + " | ".join("---" for _ in rows[0]) + " |")
    for row in rows[1:]:
        lines.append("| " + " | ".join(row) + " |")

    return "\n".join(lines)


def _process_pptx(pptx_path: Path) -> str:
    """Extract text and tables from PowerPoint .pptx files via python-pptx."""
    try:
        from pptx import Presentation
    except ImportError:
        logger.warning("python-pptx not installed, skipping %s", pptx_path)
        return ""

    try:
        prs = Presentation(str(pptx_path))
        parts: list[str] = []

        for i, slide in enumerate(prs.slides, 1):
            slide_parts: list[str] = []
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    slide_parts.append(shape.text.strip())
                if shape.has_table:
                    rows: list[list[str]] = []
                    for row in shape.table.rows:
                        cells = [cell.text.strip().replace("|", "\\|")
                                 for cell in row.cells]
                        rows.append(cells)
                    if rows:
                        max_cols = max(len(r) for r in rows)
                        for r in rows:
                            while len(r) < max_cols:
                                r.append("")
                        lines = ["| " + " | ".join(rows[0]) + " |",
                                 "| " + " | ".join("---" for _ in rows[0]) + " |"]
                        for row in rows[1:]:
                            lines.append("| " + " | ".join(row) + " |")
                        slide_parts.append("\n".join(lines))
            if slide_parts:
                parts.append(f"### Slide {i}\n\n" + "\n\n".join(slide_parts))

        return "\n\n".join(parts)

    except Exception as e:
        logger.warning("Failed to read pptx %s: %s", pptx_path, e)
        return ""


def _process_ppt_legacy(ppt_path: Path) -> str:
    """Best-effort text extraction from legacy .ppt (OLE2) files.

    Reads all OLE streams and extracts printable UTF-16LE text runs.
    Not as clean as .pptx but captures slide text content.
    """
    try:
        import olefile
    except ImportError:
        logger.warning("olefile not installed, skipping legacy .ppt %s", ppt_path)
        return ""

    try:
        ole = olefile.OleFileIO(str(ppt_path))
    except Exception as e:
        logger.warning("Cannot open OLE file %s: %s", ppt_path, e)
        return ""

    try:
        text_runs: list[str] = []
        for stream_path in ole.listdir():
            try:
                data = ole.openstream(stream_path).read()
                # PowerPoint stores text as UTF-16LE
                decoded = data.decode("utf-16-le", errors="replace")
                # Extract substantial printable runs
                runs = re.findall(r"[^\x00-\x08\x0e-\x1f]{20,}", decoded)
                text_runs.extend(runs)
            except Exception:
                continue

        # Deduplicate while preserving order
        seen: set[str] = set()
        unique: list[str] = []
        for run in text_runs:
            cleaned = run.strip()
            if cleaned and cleaned not in seen and len(cleaned) > 20:
                seen.add(cleaned)
                unique.append(cleaned)

        return "\n\n".join(unique)

    except Exception as e:
        logger.warning("Failed to parse legacy .ppt %s: %s", ppt_path, e)
        return ""
    finally:
        ole.close()


def _process_excel(excel_path: Path) -> str:
    """Convert Excel file to markdown tables."""
    try:
        import pandas as pd
    except ImportError:
        logger.warning("pandas not installed, skipping %s", excel_path)
        return ""

    try:
        engine = "openpyxl" if excel_path.suffix == ".xlsx" else "xlrd"
        xls = pd.ExcelFile(str(excel_path), engine=engine)
    except Exception as e:
        logger.warning("Cannot open Excel file %s: %s", excel_path, e)
        return ""

    parts: list[str] = []
    for sheet_name in xls.sheet_names:
        try:
            df = pd.read_excel(xls, sheet_name=sheet_name, dtype=str)
            if df.empty:
                continue
            # Limit rows for very large sheets
            if len(df) > 200:
                df = df.head(200)
            md = _dataframe_to_markdown(df)
            if md:
                parts.append(f"### Sheet: {sheet_name}\n\n{md}")
        except Exception as e:
            logger.debug("Failed to read sheet '%s' from %s: %s",
                         sheet_name, excel_path, e)

    return "\n\n".join(parts)


def _process_csv(csv_path: Path, delimiter: str = ",") -> str:
    """Convert CSV/TSV to markdown table."""
    try:
        import pandas as pd
    except ImportError:
        logger.warning("pandas not installed, skipping %s", csv_path)
        return ""

    try:
        df = pd.read_csv(str(csv_path), sep=delimiter, dtype=str, nrows=200)
        if df.empty:
            return ""
        return _dataframe_to_markdown(df)
    except Exception as e:
        logger.warning("Failed to read CSV %s: %s", csv_path, e)
        return ""


def _dataframe_to_markdown(df) -> str:
    """Convert a pandas DataFrame to a markdown table."""
    headers = list(df.columns)

    # Clean headers and values
    headers = [str(h).replace("|", "\\|").strip() for h in headers]

    lines: list[str] = []
    lines.append("| " + " | ".join(headers) + " |")
    lines.append("| " + " | ".join("---" for _ in headers) + " |")

    for _, row in df.iterrows():
        cells = [str(v).replace("|", "\\|").replace("\n", " ").strip()
                 if str(v) != "nan" else ""
                 for v in row]
        lines.append("| " + " | ".join(cells) + " |")

    return "\n".join(lines)
